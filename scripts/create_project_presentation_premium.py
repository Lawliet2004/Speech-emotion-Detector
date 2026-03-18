#!/usr/bin/env python3
from __future__ import annotations

import csv
import json
import math
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = ROOT / "presentations" / "speech_emotion_project_summary_premium.pptx"
ASSET_DIR = ROOT / "presentations" / "assets"
SPLIT_SUMMARY_PATH = ROOT / "manifests" / "split_summary.json"
AUDIT_PATH = ROOT / "reports" / "data_audit.json"
PROCESSED_MANIFEST_PATH = ROOT / "manifests" / "processed" / "processed_audio_manifest.csv"


NAVY = RGBColor(16, 24, 40)
SLATE = RGBColor(31, 41, 55)
BLUE = RGBColor(55, 118, 255)
TEAL = RGBColor(21, 173, 158)
CORAL = RGBColor(255, 117, 85)
GOLD = RGBColor(246, 189, 22)
MIST = RGBColor(243, 247, 252)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(64, 72, 88)
FAINT = RGBColor(218, 227, 240)


def load_data() -> dict:
    with SPLIT_SUMMARY_PATH.open(encoding="utf-8") as file:
        split_summary = json.load(file)
    with AUDIT_PATH.open(encoding="utf-8") as file:
        audit = json.load(file)
    with PROCESSED_MANIFEST_PATH.open(newline="", encoding="utf-8") as file:
        processed_rows = list(csv.DictReader(file))
    return {
        "split_summary": split_summary,
        "audit": audit,
        "processed_rows": processed_rows,
    }


def make_assets() -> tuple[Path, Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    cover_path = ASSET_DIR / "cover_bg.png"
    section_path = ASSET_DIR / "section_bg.png"

    if not cover_path.exists():
        img = Image.new("RGB", (1600, 900), (16, 24, 40))
        draw = ImageDraw.Draw(img, "RGBA")
        for i in range(0, 1600, 80):
            draw.ellipse((i - 50, 680, i + 220, 950), fill=(27, 66, 123, 55))
        for offset, color in [
            (0, (55, 118, 255, 170)),
            (28, (21, 173, 158, 150)),
            (56, (255, 117, 85, 120)),
        ]:
            points = []
            for x in range(0, 1601, 8):
                y = 510 + int(85 * math.sin((x / 145.0) + offset))
                points.append((x, y))
            draw.line(points, fill=color, width=5)
        draw.rounded_rectangle((1060, 110, 1480, 260), radius=28, fill=(255, 255, 255, 28))
        draw.rounded_rectangle((1085, 140, 1465, 230), radius=22, fill=(255, 255, 255, 16))
        img.save(cover_path)

    if not section_path.exists():
        img = Image.new("RGB", (1600, 900), (243, 247, 252))
        draw = ImageDraw.Draw(img, "RGBA")
        draw.rectangle((0, 0, 1600, 160), fill=(16, 24, 40, 255))
        for i in range(0, 1600, 120):
            draw.ellipse((i - 90, 640, i + 180, 930), fill=(55, 118, 255, 28))
        for x in range(0, 1601, 10):
            y = 720 + int(35 * math.sin(x / 95.0))
            draw.ellipse((x, y, x + 4, y + 4), fill=(21, 173, 158, 80))
        img.save(section_path)

    return cover_path, section_path


def set_font(run, size: int, color: RGBColor, bold: bool = False, name: str = "Aptos") -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_textbox(slide, left, top, width, height, text, size=20, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size=size, color=color, bold=bold)
    return box


def add_card(slide, left, top, width, height, fill_rgb, title, body, title_color=WHITE, body_color=WHITE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = fill_rgb
    add_textbox(slide, left + Inches(0.18), top + Inches(0.14), width - Inches(0.3), Inches(0.34), title, size=16, color=title_color, bold=True)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.48), width - Inches(0.3), height - Inches(0.58), body, size=18, color=body_color)
    return shape


def add_metric_card(slide, left, top, width, height, label, value, accent):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = FAINT
    accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.12), height)
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.color.rgb = accent
    add_textbox(slide, left + Inches(0.22), top + Inches(0.18), width - Inches(0.3), Inches(0.25), label, size=14, color=TEXT, bold=True)
    add_textbox(slide, left + Inches(0.22), top + Inches(0.48), width - Inches(0.3), Inches(0.45), value, size=24, color=NAVY, bold=True)


def add_footer(slide, page_num: int):
    add_textbox(slide, Inches(12.1), Inches(7.0), Inches(0.8), Inches(0.22), str(page_num), size=10, color=RGBColor(120, 132, 148), align=PP_ALIGN.RIGHT)


def build_presentation(data: dict) -> Presentation:
    cover_path, section_path = make_assets()
    split_summary = data["split_summary"]
    audit = data["audit"]
    processed_rows = data["processed_rows"]

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(cover_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(7.0), Inches(0.8), "Speech Emotion Detection System", size=28, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.82), Inches(1.88), Inches(5.8), Inches(1.5), "Project summary deck covering dataset work, data quality checks, preprocessing, and GPU-ready training setup.", size=18, color=RGBColor(219, 231, 247))
    add_card(slide, Inches(0.82), Inches(5.4), Inches(2.0), Inches(1.05), BLUE, "Goal", "Real-time emotion detection")
    add_card(slide, Inches(3.0), Inches(5.4), Inches(1.7), Inches(1.05), TEAL, "Platform", "Local desktop")
    add_card(slide, Inches(4.9), Inches(5.4), Inches(2.1), Inches(1.05), CORAL, "Output", "Emotion + confidence")
    add_footer(slide, 1)

    # Slide 2: Vision
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(section_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    add_textbox(slide, Inches(0.7), Inches(0.55), Inches(7.0), Inches(0.45), "Project Vision", size=26, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.7), Inches(1.35), Inches(5.6), Inches(1.4), "Build a production-minded speech emotion recognition pipeline that can take live speech, identify the dominant emotion, and support a future desktop real-time application.", size=22, color=NAVY, bold=True)
    add_card(slide, Inches(0.75), Inches(3.05), Inches(3.85), Inches(1.6), WHITE, "Why this project matters", "It combines speech processing, data quality work, reproducible ML pipelines, and real-time inference planning.", title_color=ACCENT_COLOR if False else BLUE, body_color=TEXT)
    add_card(slide, Inches(4.85), Inches(3.05), Inches(3.85), Inches(1.6), WHITE, "Current V1 scope", "Controlled desktop environment, six emotion classes, and a clean training pipeline before advanced model work.", title_color=TEAL, body_color=TEXT)
    add_card(slide, Inches(8.95), Inches(3.05), Inches(3.6), Inches(1.6), WHITE, "Delivery strategy", "Start with a solid baseline, validate results, then expand toward real-time inference and model improvement.", title_color=CORAL, body_color=TEXT)
    add_footer(slide, 2)

    # Slide 3: Dataset snapshot
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(section_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    add_textbox(slide, Inches(0.7), Inches(0.55), Inches(7.0), Inches(0.45), "Dataset Snapshot", size=24, color=WHITE, bold=True)
    add_metric_card(slide, Inches(0.75), Inches(1.25), Inches(2.35), Inches(1.1), "Raw clips", f"{split_summary['total_clips']:,}", BLUE)
    add_metric_card(slide, Inches(3.25), Inches(1.25), Inches(2.35), Inches(1.1), "Speakers", f"{split_summary['total_speakers']}", TEAL)
    add_metric_card(slide, Inches(5.75), Inches(1.25), Inches(2.35), Inches(1.1), "Processed clips", f"{len(processed_rows):,}", CORAL)
    add_metric_card(slide, Inches(8.25), Inches(1.25), Inches(2.9), Inches(1.1), "Audio spec", "Mono / 16k / 16-bit", GOLD)
    chart_data = CategoryChartData()
    chart_data.categories = [k.title() for k in split_summary["emotion_counts"].keys()]
    chart_data.add_series("Clips", list(split_summary["emotion_counts"].values()))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.85),
        Inches(2.75),
        Inches(6.5),
        Inches(3.5),
        chart_data,
    ).chart
    chart.has_legend = False
    chart.value_axis.visible = False
    chart.category_axis.tick_labels.font.size = Pt(12)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = BLUE
    chart.series[0].data_labels.show_value = True
    chart.series[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    chart.series[0].data_labels.font.size = Pt(11)
    add_card(slide, Inches(7.75), Inches(2.75), Inches(4.65), Inches(3.5), NAVY, "Key takeaway", "The dataset is already well structured for supervised learning: fixed audio format, balanced class counts, and a manageable number of speakers for speaker-disjoint evaluation.", title_color=WHITE, body_color=RGBColor(225, 232, 243))
    add_footer(slide, 3)

    # Slide 4: Workflow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_textbox(slide, Inches(0.7), Inches(0.55), Inches(7.0), Inches(0.45), "Project Workflow", size=24, color=NAVY, bold=True)
    stages = [
        ("Raw Audio", BLUE),
        ("Manifest", TEAL),
        ("Audit", CORAL),
        ("Preprocess", GOLD),
        ("Train", BLUE),
        ("Real-time", TEAL),
    ]
    left = 0.65
    for index, (label, color) in enumerate(stages):
        width = 1.9 if index < len(stages) - 1 else 1.95
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(left), Inches(2.6), Inches(width), Inches(1.15))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        add_textbox(slide, Inches(left + 0.18), Inches(2.93), Inches(width - 0.3), Inches(0.35), label, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        left += width - 0.12
    add_card(slide, Inches(0.9), Inches(4.55), Inches(3.7), Inches(1.35), WHITE, "What was implemented", "Manifest generation, data audit, duplicate handling, preprocessing, GPU training notebook, and train.py CLI.", title_color=BLUE, body_color=TEXT)
    add_card(slide, Inches(4.85), Inches(4.55), Inches(3.7), Inches(1.35), WHITE, "What is ready now", "A processed dataset and a repeatable training path that can later feed a real-time inference script.", title_color=TEAL, body_color=TEXT)
    add_card(slide, Inches(8.8), Inches(4.55), Inches(3.7), Inches(1.35), WHITE, "What comes next", "Install PyTorch locally, run the baseline training, measure performance, then move to inference and product polish.", title_color=CORAL, body_color=TEXT)
    add_footer(slide, 4)

    # Slide 5: Split strategy
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(section_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    add_textbox(slide, Inches(0.7), Inches(0.55), Inches(7.0), Inches(0.45), "Speaker-Disjoint Split Strategy", size=24, color=WHITE, bold=True)
    split_chart_data = CategoryChartData()
    split_chart_data.categories = [k.title() for k in split_summary["split_counts"].keys()]
    split_chart_data.add_series("Clips", [v["clips"] for v in split_summary["split_counts"].values()])
    split_chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.9),
        Inches(1.55),
        Inches(6.4),
        Inches(4.7),
        split_chart_data,
    ).chart
    split_chart.has_legend = False
    split_chart.value_axis.visible = False
    split_chart.category_axis.tick_labels.font.size = Pt(14)
    split_chart.series[0].format.fill.solid()
    split_chart.series[0].format.fill.fore_color.rgb = TEAL
    split_chart.series[0].data_labels.show_value = True
    split_chart.series[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    split_chart.series[0].data_labels.font.size = Pt(12)
    add_card(slide, Inches(7.8), Inches(1.65), Inches(4.45), Inches(1.2), NAVY, "Train", "5,228 clips\n64 speakers", body_color=RGBColor(220, 230, 245))
    add_card(slide, Inches(7.8), Inches(3.0), Inches(4.45), Inches(1.2), BLUE, "Validation", "1,066 clips\n13 speakers", body_color=WHITE)
    add_card(slide, Inches(7.8), Inches(4.35), Inches(4.45), Inches(1.2), CORAL, "Test", "1,147 clips\n14 speakers", body_color=WHITE)
    add_textbox(slide, Inches(7.8), Inches(5.85), Inches(4.4), Inches(0.65), "Reason for this strategy: it prevents the model from memorizing speaker traits across splits and gives a more honest estimate of generalization.", size=15, color=TEXT)
    add_footer(slide, 5)

    # Slide 6: Audit dashboard
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_textbox(slide, Inches(0.7), Inches(0.55), Inches(7.0), Inches(0.45), "Data Audit Dashboard", size=24, color=NAVY, bold=True)
    add_metric_card(slide, Inches(0.8), Inches(1.35), Inches(2.35), Inches(1.05), "Missing files", "0", TEAL)
    add_metric_card(slide, Inches(3.35), Inches(1.35), Inches(2.35), Inches(1.05), "Unreadable files", "0", TEAL)
    add_metric_card(slide, Inches(5.9), Inches(1.35), Inches(2.35), Inches(1.05), "Split leakage", "0", TEAL)
    add_metric_card(slide, Inches(8.45), Inches(1.35), Inches(2.35), Inches(1.05), "Duplicate groups", "3", CORAL)
    add_metric_card(slide, Inches(0.8), Inches(2.7), Inches(2.35), Inches(1.05), "Duration outliers", "163", GOLD)
    add_metric_card(slide, Inches(3.35), Inches(2.7), Inches(2.35), Inches(1.05), "Audit status", audit["status"].upper(), BLUE)
    add_card(slide, Inches(6.0), Inches(2.65), Inches(6.0), Inches(2.55), WHITE, "Interpretation", "The dataset is fundamentally healthy. The only notable quality issue is a very small set of exact duplicate audio pairs, which were handled conservatively during preprocessing.", title_color=NAVY, body_color=TEXT)
    add_textbox(slide, Inches(0.82), Inches(4.35), Inches(4.7), Inches(0.8), "Duration summary\nMedian 2.50 s | 95th percentile 3.47 s | Max 5.005 s", size=18, color=TEXT, bold=True)
    add_footer(slide, 6)

    # Slide 7: Preprocessing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(section_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    add_textbox(slide, Inches(0.7), Inches(0.55), Inches(7.0), Inches(0.45), "Preprocessing Pipeline", size=24, color=WHITE, bold=True)
    add_card(slide, Inches(0.82), Inches(1.45), Inches(2.1), Inches(1.2), NAVY, "Input", "Raw manifest + AudioWAV")
    add_card(slide, Inches(3.12), Inches(1.45), Inches(2.1), Inches(1.2), BLUE, "Filtering", "Exclude exact duplicate pairs")
    add_card(slide, Inches(5.42), Inches(1.45), Inches(2.1), Inches(1.2), TEAL, "Standardize", "Mono 16k PCM\nLight normalization")
    add_card(slide, Inches(7.72), Inches(1.45), Inches(2.1), Inches(1.2), CORAL, "Write outputs", "processed_audio/")
    add_card(slide, Inches(10.02), Inches(1.45), Inches(2.1), Inches(1.2), GOLD, "Manifest", "processed_audio_manifest.csv", title_color=NAVY, body_color=NAVY)
    add_metric_card(slide, Inches(1.05), Inches(3.4), Inches(2.5), Inches(1.1), "Processed clips", f"{len(processed_rows):,}", BLUE)
    add_metric_card(slide, Inches(3.95), Inches(3.4), Inches(2.5), Inches(1.1), "Excluded duplicates", "6 files", CORAL)
    add_metric_card(slide, Inches(6.85), Inches(3.4), Inches(2.5), Inches(1.1), "Failed rows", "0", TEAL)
    add_metric_card(slide, Inches(9.75), Inches(3.4), Inches(2.5), Inches(1.1), "Output status", "Ready", GOLD)
    add_textbox(slide, Inches(1.0), Inches(5.1), Inches(11.2), Inches(0.85), "Result: the project now has a clean, reproducible processed dataset and a manifest that is ready for model training.", size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 7)

    # Slide 8: Training stack
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_textbox(slide, Inches(0.7), Inches(0.55), Inches(7.0), Inches(0.45), "Training Stack", size=24, color=NAVY, bold=True)
    add_card(slide, Inches(0.9), Inches(1.4), Inches(3.45), Inches(4.85), NAVY, "Baseline design", "Processed audio -> log spectrogram -> compact 2D CNN -> emotion prediction\n\nArtifacts planned:\ncheckpoint\ntraining history\nmetrics summary\nclassification report\nconfusion matrix", title_color=WHITE, body_color=RGBColor(225, 232, 243))
    add_card(slide, Inches(4.7), Inches(1.4), Inches(3.6), Inches(2.1), WHITE, "Notebook path", "notebooks/train_baseline_gpu.ipynb", title_color=BLUE, body_color=TEXT)
    add_card(slide, Inches(4.7), Inches(3.8), Inches(3.6), Inches(2.1), WHITE, "CLI path", "train.py", title_color=TEAL, body_color=TEXT)
    add_card(slide, Inches(8.65), Inches(1.4), Inches(3.8), Inches(4.5), MIST, "Current blocker", "This sandbox does not have PyTorch installed, so the code could be written and checked structurally, but not trained here yet.\n\nNext action on the user's machine:\ninstall CUDA-enabled PyTorch and run the training entrypoint.", title_color=CORAL, body_color=TEXT)
    add_footer(slide, 8)

    # Slide 9: Project assets
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(section_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    add_textbox(slide, Inches(0.7), Inches(0.55), Inches(7.0), Inches(0.45), "What Has Been Built", size=24, color=WHITE, bold=True)
    asset_cards = [
        ("Manifest generation", "scripts/create_audio_manifest.py", BLUE),
        ("Audit pipeline", "scripts/audit_dataset.py", TEAL),
        ("Preprocessing notebook", "notebooks/preprocess_audio.ipynb", CORAL),
        ("Training notebook", "notebooks/train_baseline_gpu.ipynb", GOLD),
        ("CLI training script", "train.py", BLUE),
        ("Project handoff docs", "skill.md + context.md", TEAL),
    ]
    positions = [
        (0.85, 1.45), (4.45, 1.45), (8.05, 1.45),
        (0.85, 3.75), (4.45, 3.75), (8.05, 3.75),
    ]
    for (title, body, color), (x, y) in zip(asset_cards, positions):
        add_card(slide, Inches(x), Inches(y), Inches(3.2), Inches(1.65), WHITE, title, body, title_color=color, body_color=TEXT)
    add_footer(slide, 9)

    # Slide 10: Roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(cover_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.65), Inches(12.25), Inches(6.1))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
    overlay.fill.transparency = 0.08
    overlay.line.color.rgb = RGBColor(255, 255, 255)
    add_textbox(slide, Inches(0.9), Inches(1.0), Inches(7.0), Inches(0.5), "Next Steps", size=26, color=WHITE, bold=True)
    roadmap = [
        ("1", "Install CUDA PyTorch", BLUE),
        ("2", "Run baseline training", TEAL),
        ("3", "Review F1 and confusion matrix", CORAL),
        ("4", "Build inference path for real time use", GOLD),
    ]
    x = 0.95
    for number, label, color in roadmap:
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(2.25), Inches(0.75), Inches(0.75))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.color.rgb = color
        add_textbox(slide, Inches(x), Inches(2.42), Inches(0.75), Inches(0.2), number, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x - 0.15), Inches(3.2), Inches(2.1), Inches(0.9), label, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        if x < 9.2:
            connector = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x + 1.25), Inches(2.45), Inches(0.95), Inches(0.35))
            connector.fill.solid()
            connector.fill.fore_color.rgb = RGBColor(255, 255, 255)
            connector.line.color.rgb = RGBColor(255, 255, 255)
        x += 3.0
    add_textbox(slide, Inches(0.95), Inches(5.25), Inches(11.0), Inches(0.9), "The project now has a strong data and pipeline foundation. The next milestone is to train the first GPU baseline model and turn it into a usable emotion detection system.", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 10)

    return prs


def main() -> None:
    data = load_data()
    prs = build_presentation(data)
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Saved premium presentation to: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
