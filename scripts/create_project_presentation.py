#!/usr/bin/env python3
from __future__ import annotations

import csv
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = ROOT / "presentations" / "speech_emotion_project_summary.pptx"
SPLIT_SUMMARY_PATH = ROOT / "manifests" / "split_summary.json"
AUDIT_PATH = ROOT / "reports" / "data_audit.json"
PROCESSED_MANIFEST_PATH = ROOT / "manifests" / "processed" / "processed_audio_manifest.csv"


TITLE_COLOR = RGBColor(20, 38, 76)
ACCENT_COLOR = RGBColor(42, 99, 168)
TEXT_COLOR = RGBColor(50, 50, 50)
LIGHT_BG = RGBColor(245, 248, 252)


def load_data() -> dict:
    with SPLIT_SUMMARY_PATH.open(encoding="utf-8") as file:
        split_summary = json.load(file)
    with AUDIT_PATH.open(encoding="utf-8") as file:
        audit = json.load(file)
    with PROCESSED_MANIFEST_PATH.open(newline="", encoding="utf-8") as file:
        processed_rows = list(csv.DictReader(file))

    return {
        "split_summary": split_summary,
        "audit": audit,
        "processed_rows": processed_rows,
    }


def style_title(shape, size: int = 28) -> None:
    paragraph = shape.text_frame.paragraphs[0]
    if paragraph.runs:
        run = paragraph.runs[0]
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = TITLE_COLOR


def style_bullets(text_frame) -> None:
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(20)
            run.font.color.rgb = TEXT_COLOR


def add_header_band(slide) -> None:
    band = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.33), Inches(0.55)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT_BG
    band.line.color.rgb = LIGHT_BG


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    title_shape = slide.shapes.title
    title_shape.text = title
    style_title(title_shape, size=30)
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    for paragraph in subtitle_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = ACCENT_COLOR


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_header_band(slide)
    slide.shapes.title.text = title
    style_title(slide.shapes.title)
    body = slide.placeholders[1].text_frame
    body.clear()
    for index, bullet in enumerate(bullets):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
    style_bullets(body)


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_title: str,
    left_items: list[str],
    right_title: str,
    right_items: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_band(slide)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.5))
    title_box.text_frame.text = title
    style_title(title_box)

    for x, heading, items in [
        (0.7, left_title, left_items),
        (6.7, right_title, right_items),
    ]:
        heading_box = slide.shapes.add_textbox(Inches(x), Inches(1.2), Inches(5.2), Inches(0.4))
        heading_box.text_frame.text = heading
        for run in heading_box.text_frame.paragraphs[0].runs:
            run.font.size = Pt(22)
            run.font.bold = True
            run.font.color.rgb = ACCENT_COLOR

        body_box = slide.shapes.add_textbox(Inches(x), Inches(1.7), Inches(5.2), Inches(4.8))
        body = body_box.text_frame
        body.clear()
        for index, item in enumerate(items):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = item
            paragraph.level = 0
        style_bullets(body)


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_band(slide)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.5))
    title_box.text_frame.text = title
    style_title(title_box)

    table = slide.shapes.add_table(
        rows=len(rows) + 1,
        cols=len(headers),
        left=Inches(0.7),
        top=Inches(1.3),
        width=Inches(12.0),
        height=Inches(4.8),
    ).table

    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_COLOR
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(255, 255, 255)

    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = value
            if row_index % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(248, 250, 253)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(15)
                    run.font.color.rgb = TEXT_COLOR


def build_presentation(data: dict) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    split_summary = data["split_summary"]
    audit = data["audit"]
    processed_rows = data["processed_rows"]

    processed_count = len(processed_rows)

    add_title_slide(
        prs,
        "Speech Emotion Detection Project",
        "10-slide summary of the work completed so far | Generated on 2026-03-18",
    )

    add_bullet_slide(
        prs,
        "1. Project Goal",
        [
            "Build a production-minded speech emotion recognition system for real-time speech.",
            "Target platform for V1: local desktop.",
            "Current output contract: emotion plus confidence for each detected speech segment.",
            "Current label set: angry, disgust, fear, happy, neutral, sad.",
        ],
    )

    add_bullet_slide(
        prs,
        "2. Dataset Overview",
        [
            f"Dataset folder: AudioWAV with {split_summary['total_clips']:,} WAV clips.",
            f"Confirmed speaker count: {split_summary['total_speakers']}.",
            "Audio format confirmed as mono, 16 kHz, 16-bit PCM.",
            "Emotion distribution is close to balanced across six classes, with neutral slightly smaller.",
        ],
    )

    add_table_slide(
        prs,
        "3. Train / Val / Test Split",
        headers=["Split", "Clips", "Speakers"],
        rows=[
            [
                split.title(),
                f"{details['clips']:,}",
                f"{details['speakers']:,}",
            ]
            for split, details in split_summary["split_counts"].items()
        ],
    )

    add_bullet_slide(
        prs,
        "4. Data Organization Work",
        [
            "Created a manifest-driven workflow instead of moving raw audio into split folders.",
            "Added scripts/create_audio_manifest.py to parse filenames and create speaker-disjoint splits.",
            "Saved artifacts: audio_manifest.csv, train_manifest.csv, val_manifest.csv, test_manifest.csv, and split_summary.json.",
            "This makes the experiment setup reproducible and safer to iterate on later.",
        ],
    )

    add_bullet_slide(
        prs,
        "5. Data Audit Findings",
        [
            f"Audit status: {audit['status'].upper()}.",
            "No missing files, unreadable WAV files, manifest/header mismatches, or split leakage were found.",
            f"Main review item: {audit['issues'][0] if audit['issues'] else 'none'}.",
            "Duration outliers were identified statistically, but they are not treated as automatic corruption.",
        ],
    )

    add_bullet_slide(
        prs,
        "6. Preprocessing Pipeline",
        [
            "Rewrote preprocessing to use only standard-library WAV handling plus NumPy/Pandas.",
            "Excluded both files from each duplicate-audio pair as a conservative label-quality policy.",
            "Standardized each kept clip and wrote outputs to data/processed_audio/.",
            "Saved a clean processed manifest for downstream model training.",
        ],
    )

    add_two_column_slide(
        prs,
        "7. Preprocessing Outputs",
        "Processed Dataset",
        [
            f"Processed WAV files written: {processed_count:,}",
            "Processed manifest rows: 7,435",
            "Failed preprocessing rows: 0",
            "Output manifest: manifests/processed/processed_audio_manifest.csv",
        ],
        "Why This Matters",
        [
            "Training now uses a consistent and validated source of truth.",
            "The data pipeline can be rerun without touching raw AudioWAV files.",
            "This reduces future training and inference inconsistencies.",
        ],
    )

    add_bullet_slide(
        prs,
        "8. Model Training Setup",
        [
            "Created a GPU-ready baseline notebook: notebooks/train_baseline_gpu.ipynb.",
            "Created a simple CLI training entrypoint: train.py.",
            "Baseline model: compact log-spectrogram CNN with mixed precision on CUDA when available.",
            "Training artifacts are designed to save checkpoints, history, metrics summary, report, and confusion matrix.",
        ],
    )

    add_bullet_slide(
        prs,
        "9. Current Status and Next Steps",
        [
            "Completed so far: dataset organization, audit, preprocessing, processed manifest generation, and baseline training code.",
            "Current blocker in this environment: PyTorch is not installed, so training could not be executed here yet.",
            "Next step: install CUDA-enabled PyTorch locally and run train.py or the training notebook.",
            "After baseline training: evaluate results, then add inference code for real-time emotion prediction.",
        ],
    )

    return prs


def main() -> None:
    data = load_data()
    prs = build_presentation(data)
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Saved presentation to: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
